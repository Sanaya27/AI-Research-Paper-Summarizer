import os
import google.generativeai as genai
import streamlit as st

from pptx import Presentation
from pptx.util import Inches, Pt

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter


# ============================================================
#  GEMINI API KEY HANDLING (LOCAL + DEPLOYMENT SAFE)
# ============================================================

# 1) Try Streamlit Cloud secrets first
if "GEMINI_API_KEY" in st.secrets:
    GEMINI_API_KEY = st.secrets["GEMINI_API_KEY"]

# 2) Fall back to system environment variable for local testing
else:
    GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:
    raise ValueError("❌ Gemini API Key missing. Add it to st.secrets or .env environment variable.")

genai.configure(api_key=GEMINI_API_KEY)

# Main model
model = genai.GenerativeModel("gemini-2.0-flash")

# ============================================================
# 1) SUMMARY GENERATOR
# ============================================================

def summarize_text(text, length="Medium", style="Academic"):

    # summary length instruction
    if length == "Short":
        length_instr = "Summarize the text briefly in 3–4 concise sentences."
    elif length == "Long":
        length_instr = "Provide a detailed multi-paragraph summary covering all major points, methods, findings and conclusions."
    else:
        length_instr = "Summarize the text in a balanced 2 paragraphs."

    # writing style
    if style == "Academic":
        style_instr = "Use an academic tone."
    elif style == "Simple English":
        style_instr = "Explain in simple, easy-to-understand English."
    elif style == "Bullet Points":
        style_instr = "Provide the summary as bullet points."
    elif style == "Technical":
        style_instr = "Use precise technical language."
    else:
        style_instr = "Write in a smooth, narrative tone."

    prompt = f"""
    {length_instr}
    {style_instr}

    Text to summarize:
    {text}
    """

    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"❌ Gemini API Error: {str(e)}"


# ============================================================
# 2) TITLE EXTRACTION
# ============================================================

def extract_title(text):
    prompt = f"""
    Extract ONLY the best possible title of this research paper.
    Respond with ONLY the title.

    Text:
    {text}
    """
    try:
        return model.generate_content(prompt).text.strip()
    except Exception as e:
        return f"❌ Gemini API Error: {str(e)}"


# ============================================================
# 3) KEYWORDS EXTRACTION
# ============================================================

def extract_keywords(text):
    prompt = f"""
    Extract the 5–10 most important keywords.
    Return them as a comma-separated list ONLY.

    Text:
    {text}
    """
    try:
        return model.generate_content(prompt).text.strip()
    except Exception as e:
        return f"❌ Gemini API Error: {str(e)}"


# ============================================================
# 4) PLAGIARISM CHECK
# ============================================================

def check_plagiarism(text):
    prompt = f"""
    You are an AI plagiarism and originality detector.

    Analyze:
    - Estimated originality % 
    - Whether text is AI-generated / copied / human-written
    - Suspicious patterns
    - Short explanation

    Do NOT rewrite the text.

    Text:
    {text}
    """

    try:
        return model.generate_content(prompt).text.strip()
    except Exception as e:
        return f"❌ Gemini API Error: {str(e)}"


# ============================================================
# 5) SEMANTIC SEARCH
# ============================================================

def semantic_search(query, text):
    prompt = f"""
    You are an AI doing semantic search inside a research paper.

    Task:
    - Understand question meaning
    - Find relevant portion of TEXT
    - Give a clear answer
    - If not found → say: "The document does not contain this information."

    Question:
    {query}

    Research Paper:
    {text}
    """

    try:
        return model.generate_content(prompt).text.strip()
    except Exception as e:
        return f"❌ Gemini API Error: {str(e)}"


# ============================================================
# 6) AUTO-GENERATED PPT
# ============================================================

def generate_ppt(text):

    prompt = f"""
    Convert this research paper into structured slide information.

    Extract EXACTLY:

    Title: <title>
    Authors: <authors or "Not specified">

    Problem:
    <2–4 lines>

    Objectives:
    <2–4 lines>

    Methodology:
    <3–6 lines>

    Results:
    <3–6 lines>

    Conclusion:
    <2–4 lines>

    Keywords: <comma separated keywords>

    Text:
    {text}
    """

    try:
        outline = model.generate_content(prompt).text.strip()
    except Exception as e:
        return f"❌ PPT Generation Error (LLM Step): {str(e)}"

    # --- Parse response
    sections = {
        "title": "Untitled Presentation",
        "authors": "Not specified",
        "problem": "",
        "objectives": "",
        "methodology": "",
        "results": "",
        "conclusion": "",
        "keywords": "",
    }

    current = None
    header_map = {
        "title:": "title",
        "authors:": "authors",
        "problem:": "problem",
        "objectives:": "objectives",
        "methodology:": "methodology",
        "results:": "results",
        "conclusion:": "conclusion",
        "keywords:": "keywords",
    }

    for line in outline.splitlines():
        stripped = line.strip()
        lower = stripped.lower()

        # header detected
        for h, key in header_map.items():
            if lower.startswith(h):
                current = key
                after = stripped.split(":", 1)[1].strip()
                if after:
                    sections[key] = after
                break
        else:
            # continuation lines
            if current:
                sections[current] += "\n" + stripped

    # --- Build PPT
    try:
        prs = Presentation()

        # title slide
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sections["title"]
        slide.placeholders[1].text = f"Authors: {sections['authors']}"

        def add_slide(title, content):
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            frame = slide.shapes.placeholders[1].text_frame
            frame.clear()

            for line in content.split("\n"):
                if line.strip():
                    p = frame.add_paragraph()
                    p.text = line.strip()
                    p.level = 0

        add_slide("Problem Statement", sections["problem"])
        add_slide("Objectives", sections["objectives"])
        add_slide("Methodology", sections["methodology"])
        add_slide("Results / Findings", sections["results"])
        add_slide("Conclusion", sections["conclusion"])
        add_slide("Keywords", sections["keywords"])

        out = "generated_presentation.pptx"
        prs.save(out)
        return out

    except Exception as e:
        return f"❌ PPT Generation Error (PPT step): {str(e)}"


# ============================================================
# 7) EXTRACT ALGORITHMS + EQUATIONS
# ============================================================

def extract_algorithms_equations(text):

    prompt = f"""
    Extract from the text:

    Equations (Math, LaTeX, symbolic)
    Algorithms (Algorithm 1, pseudocode, steps)

    Return format:

    Equations:
    - eq1
    - eq2

    Algorithms:
    - step1
    - step2

    TEXT:
    {text}
    """

    try:
        return model.generate_content(prompt).text.strip()
    except Exception as e:
        return f"❌ Equation Extraction Error: {str(e)}"


# ============================================================
# 8) PDF RESEARCH NOTES
# ============================================================

def generate_research_notes_pdf(title, keywords, summary, plagiarism_report, algorithms_equations):

    try:
        filename = "Research_Notes.pdf"

        doc = SimpleDocTemplate(
            filename,
            pagesize=letter,
            rightMargin=40, leftMargin=40,
            topMargin=50, bottomMargin=50
        )

        styles = getSampleStyleSheet()
        story = []

        def add_section(header, text):
            story.append(Paragraph(f"<b>{header}</b>", styles['Heading2']))
            story.append(Paragraph(text.replace("\n", "<br/>"), styles['BodyText']))
            story.append(Spacer(1, 12))

        # Title
        story.append(Paragraph(f"<b>{title}</b>", styles['Title']))
        story.append(Spacer(1, 20))

        add_section("Keywords:", keywords)
        add_section("Summary:", summary)
        add_section("Algorithms & Equations:", algorithms_equations)
        add_section("Plagiarism Report:", plagiarism_report)

        doc.build(story)

        return filename

    except Exception as e:
        return f"❌ PDF Generation Error: {str(e)}"
